# -*- coding: utf-8 -*-
"""
校园外卖两段式配送系统 · 答辩汇报 PPT 精修 v5.0
基于 python-pptx 从零生成，统一深蓝主题，36-38页
"""
import sys, io
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy
import os

# ============================================================
# Design Constants
# ============================================================
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
DARK_BLUE = RGBColor(0x2F, 0x54, 0x96)
LIGHT_BLUE = RGBColor(0x4A, 0x7F, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xB0, 0x50)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
RED = RGBColor(0xCC, 0x33, 0x33)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_BG = RGBColor(0x1A, 0x23, 0x33)
TITLE_BAR_H = Inches(0.9)
FONT = '微软雅黑'  # 微软雅黑
# Use direct Chinese for font name
FONT_CN = '微软雅黑'
MONO_FONT = 'Consolas'

IMG_DIR = os.path.join(os.path.dirname(__file__), 'images', 'screenshots')
ER_PATH = os.path.join(os.path.dirname(__file__), 'images', 'er_diagram.png')


def add_title_bar(slide, title_text, page_num=None):
    """Add standard blue title bar at top"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, TITLE_BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.LEFT
    # vertically center
    tf.paragraphs[0].space_before = Pt(6)

    if page_num is not None:
        add_page_number(slide, page_num)


def add_page_number(slide, num):
    """Add page number at bottom right"""
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = FONT_CN
    p.alignment = PP_ALIGN.RIGHT


def add_section_header(prs, number, title, subtitle):
    """Dark section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full dark background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    # Number
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(1.5), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(60)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_CN

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(10), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_CN

    # Subtitle
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(4.0), Inches(10), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_TEXT
    p.font.name = FONT_CN


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=DARK_TEXT, bold=False, font_name=FONT_CN, align=PP_ALIGN.LEFT):
    """Helper to add a text box"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=DARK_TEXT, font_name=FONT_CN, line_spacing=1.3):
    """Add text box with multiple paragraphs"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
    return tf


def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    """Add a code block with dark background"""
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x26, 0x2B, 0x33)
    bg.line.color.rgb = RGBColor(0x3A, 0x3F, 0x47)
    bg.line.width = Pt(1)
    # Text
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.1), Inches(width - 0.3), Inches(height - 0.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = MONO_FONT
        p.space_after = Pt(1)
        # Color: keywords in light blue, rest in light gray
        if line.strip().startswith('--') or line.strip().startswith('#'):
            p.font.color.rgb = RGBColor(0x6A, 0x99, 0x5A)
        else:
            p.font.color.rgb = RGBColor(0xDC, 0xDC, 0xDC)
    return tf


def add_info_box(slide, left, top, width, height, title, lines, box_color=DARK_BLUE):
    """Add a colored info/explanation box"""
    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(
        min(box_color[0] + 30, 255) if isinstance(box_color, tuple) else
        min(int(box_color[1:3], 16) + 30, 255) if isinstance(box_color, str) else 240,
        min(int(box_color[3:5], 16) + 30, 255) if isinstance(box_color, str) else 240,
        min(int(box_color[5:7], 16) + 30, 255) if isinstance(box_color, str) else 245
    )
    bg.fill.fore_color.rgb = RGBColor(0xE3, 0xEE, 0xF7)  # Light blue bg
    bg.line.color.rgb = DARK_BLUE
    bg.line.width = Pt(1.5)

    # Title
    txBox = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.08), Inches(width - 0.3), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    p.font.name = FONT_CN

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.4), Inches(width - 0.3), Inches(height - 0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_CN
        p.space_after = Pt(2)


def add_card(slide, left, top, width, height, title, content_lines, accent_color=DARK_BLUE):
    """Add a card with accent border"""
    # White card bg
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(0.5)

    # Top accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Title
    add_textbox(slide, left + 0.15, top + 0.12, width - 0.3, 0.35,
                title, font_size=14, color=accent_color, bold=True)

    # Content
    add_multiline_textbox(slide, left + 0.15, top + 0.5, width - 0.3, height - 0.55,
                          content_lines, font_size=11, color=DARK_TEXT)


def add_table_simple(slide, left, top, col_widths, headers, rows, font_size=11):
    """Add a simple table"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                         Inches(left), Inches(top),
                                         Inches(total_w), Inches(0.35 * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_CN
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size - 1)
                p.font.color.rgb = DARK_TEXT
                p.font.name = FONT_CN
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
    return table_shape


def insert_image_safe(slide, path, left, top, width=None, height=None):
    """Insert image if file exists"""
    if os.path.exists(path):
        kwargs = {'left': Inches(left), 'top': Inches(top)}
        if width:
            kwargs['width'] = Inches(width)
        if height:
            kwargs['height'] = Inches(height)
        return slide.shapes.add_picture(path, **kwargs)
    else:
        print(f"  [WARN] Image not found: {path}")
        return None


# ============================================================
# BUILD PPT
# ============================================================

def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    pg = [0]  # page counter (mutable)

    def page(add=1):
        pg[0] += add
        return pg[0]

    # ===== SLIDE 1: COVER =====
    slide = prs.slides.add_slide(BLANK)
    # Blue gradient-ish bg
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = DARK_BLUE
    stripe.line.fill.background()

    add_textbox(slide, 1.0, 1.2, 11, 0.6, '数据库系统课程项目 · 答辩汇报', 16, GRAY_TEXT)
    add_textbox(slide, 1.0, 1.9, 11, 1.2,
                '校园外卖两段式配送\n数据库系统', 42, WHITE, True)
    add_textbox(slide, 1.0, 3.5, 11, 1.0,
                'MySQL 8.0 + Flask + ECharts + DeepSeek AI  |  7表 4索引 7触发器 4存储过程 2视图',
                14, GRAY_TEXT)
    add_textbox(slide, 1.0, 5.5, 11, 0.5, 'Campus Delivery Two-Stage Distribution Database System', 12, GRAY_TEXT)
    add_textbox(slide, 1.0, 6.3, 11, 0.5, '2026年6月', 14, GRAY_TEXT)
    page()

    # ===== SLIDE 2: OUTLINE =====
    add_section_header(prs, '目 录', 'CONTENTS',
                       '项目背景 | 系统设计 | ER图与7表 | 六态状态机 | 索引策略 | 七重触发器 | 存储过程与事务 | 视图 | 大屏展示 | 创新总结')
    page()

    # ===== SLIDE 3: SECTION 01 - Project Background =====
    add_section_header(prs, '01', '项目背景与需求分析',
                       '校园外卖市场规模 · 传统配送痛点 · 两段式模型解决方案')
    page()

    # ===== SLIDE 4: Market & Pain Points =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '校园外卖市场与核心痛点', page())

    # Market stats - 3 cards
    cards_data = [
        ('4,500亿+', '中国外卖市场规模(2025)', '年复合增长率 22.4%', GREEN),
        ('70%+', '高校外卖渗透率', '大学生是外卖消费主力军', DARK_BLUE),
        ('120万单/日', '头部高校日均外卖量', '高峰期配送压力巨大', ORANGE),
    ]
    for i, (num, title, sub, clr) in enumerate(cards_data):
        x = 0.8 + i * 4.0
        # Card bg
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.4), Inches(3.5), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        card.line.width = Pt(0.5)
        add_textbox(slide, x + 0.2, 1.5, 3.1, 0.6, num, 32, clr, True)
        add_textbox(slide, x + 0.2, 2.2, 3.1, 0.4, title, 14, DARK_TEXT, True)
        add_textbox(slide, x + 0.2, 2.6, 3.1, 0.4, sub, 11, GRAY_TEXT)

    # Pain points
    add_textbox(slide, 0.8, 3.6, 11, 0.5, '传统校园外卖配送的三大痛点', 18, DARK_BLUE, True)
    pains = [
        '[1] 校门禁入 — 校外骑手被保安拦下，无法进入宿舍区，学生需步行数百米取餐',
        '[2] 高峰拥堵 — 午晚餐高峰期 30+ 骑手同时在校门口等学生，场面混乱，配送时效无法保证',
        '[3] 丢失错拿 — 外卖堆放在地上，无管理无追溯，日均丢失率 3%-5%，责任无法界定',
    ]
    add_multiline_textbox(slide, 0.8, 4.2, 11.5, 2.0, pains, 14, DARK_TEXT, line_spacing=1.5)

    # ===== SLIDE 5: Two-Stage Model =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '两段式配送模型：破解传统痛点', page())

    add_textbox(slide, 0.8, 1.2, 11, 0.5, '核心思路：将一条配送链拆分为两段，引入宿舍楼寄存柜作为中转枢纽', 16, DARK_BLUE, True)

    # Stage 1 card
    add_card(slide, 0.5, 1.9, 5.8, 2.8,
             '第一段：干线配送（商家 → 寄存点）',
             ['骑手类型：Stage1_Trunk（干线骑手）',
              '流程：用户下单 → 付款 → 指派干线骑手 → 商家取餐 → 送达宿舍寄存柜',
              '关键约束：柜子容量有限（物理格子数），满了不能入库',
              '关键技术：FOR UPDATE行级锁防超卖 + chk_capacity容量约束',
              '送达后：订单状态变为 Arrived_At_Point，干线骑手自动释放为 Idle'],
             DARK_BLUE)

    # Stage 2 card
    add_card(slide, 6.8, 1.9, 5.8, 2.8,
             '第二段：楼栋配送（寄存点 → 寝室）',
             ['骑手类型：Stage2_Floor（楼栋骑手）',
              '流程：寄存柜包裹就绪 → 指派楼栋骑手 → 从柜中取件 → 配送至学生寝室门口',
              '时效优势：楼栋骑手熟悉本楼布局，单次取多件批量配送',
              '关键技术：骑手类型ENUM约束 + 状态自动管理触发器',
              '送达后：订单状态变为 Completed，包裹从柜中扣除，楼栋骑手释放'],
             GREEN)

    # Bottom note
    add_textbox(slide, 0.5, 5.1, 12, 0.8,
                '核心价值：校外骑手只需送到寄存点（不进入宿舍区），楼栋骑手负责最后100米 — 安全、高效、可追溯',
                13, DARK_TEXT)

    # ===== SLIDE 6: State Machine Overview =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '六态订单状态机：全生命周期追踪', page())

    # Flow diagram using text boxes
    states = [
        ('Paid', '已支付\n待指派', RGBColor(0x60, 0x7D, 0x8B)),
        ('Stage1\nAssigned', '干线骑手\n已接单', RGBColor(0x45, 0x60, 0x97)),
        ('Arrived\nAt_Point', '已送达\n寄存点', RGBColor(0x2F, 0x54, 0x96)),
        ('Stage2\nAssigned', '楼栋骑手\n已接单', RGBColor(0x00, 0xB0, 0x50)),
        ('Completed', '已送达\n完成', RGBColor(0x2E, 0x7D, 0x32)),
        ('Cancelled', '已取消', RED),
    ]
    for i, (name, desc, clr) in enumerate(states):
        x = 0.3 + i * 2.1
        # State box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(1.8), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = clr
        box.line.fill.background()
        box.adjustments[0] = 0.15
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER

        add_textbox(slide, x, 2.8, 1.8, 0.6, desc, 10, GRAY_TEXT, align=PP_ALIGN.CENTER)

        # Arrow
        if i < len(states) - 2:
            add_textbox(slide, x + 1.8, 1.8, 0.3, 0.5, '▶', 16, GRAY_TEXT)

    # Explanation
    add_textbox(slide, 0.8, 3.5, 11, 0.5, '订单状态流转说明', 16, DARK_BLUE, True)
    explanations = [
        'Paid → Stage1_Assigned: 系统指派干线骑手取件',
        'Stage1_Assigned → Arrived_At_Point: 干线骑手送达寄存柜(调用 sp_arrive_at_pickup_point)',
        'Arrived_At_Point → Stage2_Assigned: 系统指派楼栋骑手从柜中取件',
        'Stage2_Assigned → Completed: 楼栋骑手送达学生(调用 sp_stage2_deliver)',
        '任意非终态 → Cancelled: 订单取消(调用 sp_cancel_order), 退款+恢复库存+释放骑手',
    ]
    add_multiline_textbox(slide, 0.8, 4.1, 11, 2.5, explanations, 13, DARK_TEXT, line_spacing=1.5)

    # Audit timestamps
    add_textbox(slide, 0.8, 6.0, 11, 0.4,
                '全链路审计: created_at → stage1_completed_at → stage2_completed_at (TIMESTAMP精准追踪)',
                12, GRAY_TEXT)

    page()

    # ===== SLIDE 7: SECTION 02 - System Design =====
    add_section_header(prs, '02', '系统设计与数据库架构',
                       '三层架构 · 技术栈 · ER图 · 7张核心表(3NF) · 4个索引')
    page()

    # ===== SLIDE 8: System Architecture =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '系统总体架构：三层分离', page())

    layers = [
        ('展现层', 'Flask Web + ECharts 数据可视化大屏\nlocalhost:5000 实时监控\n3个Tab: 总览 | 订单管理 | 爆仓预警', DARK_BLUE),
        ('业务逻辑层', 'Python Flask REST API\n4个存储过程(下单/入库/配送/取消)\nAI Text-to-SQL(DeepSeek) 自然语言查数据', GREEN),
        ('数据持久层', 'MySQL 8.0 InnoDB\n7张表 | 7个触发器 | 2个视图 | 4个索引\nFOR UPDATE行级锁 | CHECK约束 | ENUM类型 | FOREIGN KEY级联', ORANGE),
    ]
    for i, (name, desc, clr) in enumerate(layers):
        y = 1.3 + i * 1.9
        # Layer box
        lbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(11.5), Inches(1.5))
        lbox.fill.solid()
        lbox.fill.fore_color.rgb = WHITE
        lbox.line.color.rgb = clr
        lbox.line.width = Pt(2)

        # Name badge
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(1.6), Inches(1.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = clr
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER

        add_multiline_textbox(slide, 2.7, y + 0.15, 9, 1.2,
                              desc.split('\n'), 13, DARK_TEXT, line_spacing=1.4)

    # ===== SLIDE 9: Tech Stack =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '技术栈全图', page())

    tech_categories = [
        ('数据库引擎', ['MySQL 8.0 (InnoDB)', 'utf8mb4 字符集', 'B-Tree 索引', 'FOR UPDATE 行级锁']),
        ('后端框架', ['Python 3.11', 'Flask 3.x', 'PyMySQL (原生SQL)', 'Flask-CORS 跨域']),
        ('前端可视化', ['ECharts 5.x', '响应式仪表盘', '3个功能Tab', '实时数据刷新']),
        ('AI 集成', ['DeepSeek API', 'Text-to-SQL', '自然语言查询', 'Schema 注入提示词']),
        ('部署工具', ['Cloudflare Tunnel', '公网访问隧道', 'Git 版本控制', '.env 配置管理']),
    ]
    for i, (cat, items) in enumerate(tech_categories):
        x = 0.3 + (i % 3) * 4.2
        y = 1.2 + (i // 3) * 2.8
        add_card(slide, x, y, 3.8, 2.4, cat, items,
                 [DARK_BLUE, GREEN, ORANGE, RGBColor(0x7B, 0x1F, 0xA2), RED][i])

    page()

    # ===== SLIDE 10: ER Diagram =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '数据库 ER 图：7张核心表 · 严格3NF', page())
    if os.path.exists(ER_PATH):
        insert_image_safe(slide, ER_PATH, 0.5, 1.2, 12.3, 5.8)
    else:
        add_textbox(slide, 1, 3, 10, 1, '[ER 图文件未找到，请检查 images/er_diagram.png]', 18, RED)
    page()

    # ===== SLIDE 11: 7 Tables Overview =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '7张核心表概览', page())

    add_table_simple(slide, 0.5, 1.3,
                     [1.5, 1.8, 1.2, 7.0],
                     ['表名', '中文名', '记录数', '核心字段与约束'],
                     [
                         ['users', '学生用户表', '5', 'user_id PK, username, phone UNIQUE, dorm_building, balance CHECK(>=0)'],
                         ['merchants', '商家表', '3', 'merchant_id PK, merchant_name, phone, address, rating CHECK(1..5)'],
                         ['dishes', '菜品表', '6', 'dish_id PK, merchant_id FK, price CHECK(>0), stock CHECK(>=0), status(0/1)'],
                         ['pickup_points', '寄存点表', '12', 'point_id PK, capacity, current_packages, chk_capacity CHECK(current<=capacity)'],
                         ['riders', '骑手表', '15', 'rider_id PK, rider_type ENUM(Stage1_Trunk/Stage2_Floor), status ENUM(Idle/Delivering/Offline)'],
                         ['orders', '订单表', '~5000', 'order_id PK, 5个FK, order_status ENUM(6态), stage1_rider_id, stage2_rider_id, 3个时间戳'],
                         ['order_items', '订单明细表', '~5000', 'item_id PK, order_id FK(CASCADE), dish_id FK, quantity, price_at_order'],
                     ], font_size=9)

    add_textbox(slide, 0.5, 5.8, 11, 0.5,
                '设计规范：严格3NF · 完备外键级联 · CHECK约束 · ENUM类型约束 · InnoDB事务引擎 · B-Tree索引优化',
                12, GRAY_TEXT)

    # ===== SLIDE 12: Table Details (Pickup Points + Riders) =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '关键表详解：pickup_points 与 riders', page())

    # Pickup points table with SQL
    add_textbox(slide, 0.5, 1.15, 6, 0.4, 'pickup_points — 宿舍寄存中转点表（两段式配送核心）', 15, DARK_BLUE, True)
    add_code_block(slide, 0.5, 1.6, 6.0, 2.2, [
        'CREATE TABLE pickup_points (',
        '  point_id INT AUTO_INCREMENT PRIMARY KEY,',
        '  point_name VARCHAR(50) NOT NULL,        -- 寄存点名称',
        '  location VARCHAR(200) NOT NULL,          -- 物理位置',
        '  capacity INT NOT NULL,                   -- 最大格子数',
        '  current_packages INT DEFAULT 0',
        '    CHECK (current_packages >= 0),         -- 在库件数',
        '  -- 物理容积硬约束: 不允许超容写入',
        '  CONSTRAINT chk_capacity',
        '    CHECK (current_packages <= capacity)',
        ');',
    ], 10)
    add_info_box(slide, 0.5, 4.0, 6.0, 2.6, '设计要点',
                 ['chk_capacity 是物理约束最后防线 — 80个格子绝不允许81个包裹',
                  'current_packages 由存储过程维护(+1入库/-1出库)，不直接修改',
                  '12个寄存点覆盖全校宿舍区，容量从50到120格不等',
                  '与 orders 表通过 pickup_point_id FK 关联'])

    # Riders table
    add_textbox(slide, 6.8, 1.15, 6, 0.4, 'riders — 两段式特种骑手表', 15, DARK_BLUE, True)
    add_code_block(slide, 6.8, 1.6, 5.8, 1.8, [
        'CREATE TABLE riders (',
        '  rider_id INT AUTO_INCREMENT PRIMARY KEY,',
        '  rider_name VARCHAR(50) NOT NULL,',
        '  rider_type ENUM(',
        "    'Stage1_Trunk',    -- 干线: 商家->寄存点",
        "    'Stage2_Floor'     -- 楼栋: 寄存点->寝室",
        '  ) NOT NULL,',
        '  status ENUM(',
        "    'Idle','Delivering','Offline'",
        '  ) DEFAULT \'Idle\',',
        ');',
    ], 10)
    add_info_box(slide, 6.8, 3.6, 5.8, 3.0, '设计要点',
                 ['ENUM 类型在数据库层强制约束骑手分工，Stage1_Trunk不能做楼栋配送',
                  'status 由触发器自动管理 — 分配订单→Delivering, 完成/取消→Idle',
                  '15个骑手: 8个干线(Stage1_Trunk) + 7个楼栋(Stage2_Floor)',
                  '外键: orders.stage1_rider_id→riders, orders.stage2_rider_id→riders',
                  '两个骑手同一订单互不冲突，各自独立管理'])

    page()

    # ===== SLIDE 13: Table Details (Orders) =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '核心表详解：orders 订单表（六态状态机引擎）', page())

    add_code_block(slide, 0.5, 1.2, 7.5, 4.5, [
        'CREATE TABLE orders (',
        '  order_id INT AUTO_INCREMENT PRIMARY KEY,    -- 订单流水号',
        '  user_id INT NOT NULL,            -- FK → users         下单学生',
        '  merchant_id INT NOT NULL,        -- FK → merchants      下单商家',
        '  pickup_point_id INT NOT NULL,    -- FK → pickup_points  指派寄存点',
        '  total_amount DECIMAL(10,2),      -- 订单总金额(=菜品单价x数量)',
        '  -- ││ 六态状态机 ││',
        "  order_status ENUM('Paid','Stage1_Assigned',",
        "    'Arrived_At_Point','Stage2_Assigned',",
        "    'Completed','Cancelled') DEFAULT 'Paid',",
        '  -- 双骑手绑定',
        '  stage1_rider_id INT DEFAULT NULL,  -- FK → riders  干线骑手',
        '  stage2_rider_id INT DEFAULT NULL,  -- FK → riders  楼栋骑手',
        '  -- 全链路时间审计戳',
        '  created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,',
        '  stage1_completed_at TIMESTAMP NULL,   -- 干线送达时间',
        '  stage2_completed_at TIMESTAMP NULL,   -- 最终送达时间',
        '  -- 5个外键约束 + 2个骑手FK',
        '  FOREIGN KEY (user_id) REFERENCES users(user_id),',
        '  FOREIGN KEY (merchant_id) REFERENCES merchants(merchant_id),',
        '  FOREIGN KEY (pickup_point_id) REFERENCES pickup_points(point_id),',
        '  FOREIGN KEY (stage1_rider_id) REFERENCES riders(rider_id),',
        '  FOREIGN KEY (stage2_rider_id) REFERENCES riders(rider_id),',
        ');',
    ], 10)

    # Right side: key design points
    add_info_box(slide, 8.3, 1.2, 4.5, 1.6, '六态 ENUM 状态机',
                 ['Paid → Stage1_Assigned → Arrived_At_Point',
                  '→ Stage2_Assigned → Completed',
                  'Cancelled 可从任意非终态进入',
                  '状态驱动配送流程，杜绝脏数据'])

    add_info_box(slide, 8.3, 3.0, 4.5, 1.3, '双骑手复合绑定',
                 ['stage1_rider_id: 干线骑手(商家→寄存点)',
                  'stage2_rider_id: 楼栋骑手(寄存点→寝室)',
                  '两个骑手状态独立，由触发器自动管理'])

    add_info_box(slide, 8.3, 4.5, 4.5, 1.8, '三时间戳审计链',
                 ['created_at: 下单时间',
                  'stage1_completed_at: 干线送达柜时间',
                  'stage2_completed_at: 最终送达学生时间',
                  '→ 全链路可追溯，配送时效可量化'])

    page()

    # ===== SLIDE 14: SECTION 03 - Indexes =====
    add_section_header(prs, '03', '索引策略与性能优化',
                       '4个B-Tree索引 · 覆盖查询加速 · 复合索引设计 · 视图关联优化')
    page()

    # ===== SLIDE 15: Indexes Detailed =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '4个B-Tree索引：加速大屏查询与容量检查', page())

    add_table_simple(slide, 0.5, 1.3,
                     [1.8, 1.6, 1.6, 7.0],
                     ['索引名', '表', '索引列', '加速场景'],
                     [
                         ['idx_orders_status', 'orders', 'order_status', '大屏 Dashboard 按状态筛选(待配送/配送中/已完成) — 最频繁查询'],
                         ['idx_orders_created', 'orders', 'created_at', '历史订单时间范围查询(今日/近7天/本月) — 支撑趋势图'],
                         ['idx_dishes_merchant', 'dishes', 'merchant_id, status', '前端点餐页加载某商家上架菜品 — 复合索引覆盖查询'],
                         ['idx_orders_point_status',
                          'orders',
                          'pickup_point_id,\norder_status',
                          '寄存点容量检查 + vw_pickup_point_analytics视图关联查询\n→ sp_create_order容量预检 + 爆仓预警视图'],
                     ], font_size=10)

    add_textbox(slide, 0.5, 4.2, 11, 0.5, '索引设计原则', 16, DARK_BLUE, True)
    principles = [
        '[1] 高频查询列优先建索引 — order_status 是大屏 Tab 切换的核心过滤条件',
        '[2] 复合索引最左前缀 — idx_dishes_merchant(merchant_id, status) 同时加速"按商家查菜品"和"按商家+上架状态查"',
        '[3] 新索引 idx_orders_point_status 解决 VIEW LEFT JOIN 全表扫描问题 — 覆盖 pickup_point_id + order_status 组合查询',
        '[4] 外键列自动建索引(InnoDB) — user_id, merchant_id, pickup_point_id, stage1_rider_id, stage2_rider_id 均有B-Tree索引',
    ]
    add_multiline_textbox(slide, 0.5, 4.8, 11.5, 2.0, principles, 13, DARK_TEXT, line_spacing=1.4)

    page()

    # ===== SLIDE 16: SECTION 04 - Triggers =====
    add_section_header(prs, '04', '七重触发器防护盾',
                       'FOR UPDATE行级锁 · 库存防超卖 · 骑手类型约束 · 状态自动管理 · 容量预检')
    page()

    # ===== SLIDE 17: FOR UPDATE Deep Dive =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, 'FOR UPDATE 行级锁：高并发防超卖核心机制', page())

    # Left: Race condition
    add_textbox(slide, 0.5, 1.15, 5.5, 0.4, '竞态条件（无锁情况）', 16, RED, True)
    add_code_block(slide, 0.5, 1.6, 5.5, 2.0, [
        '-- 时刻 T1: Session A 查询库存',
        'SELECT stock FROM dishes',
        "  WHERE dish_id=1;          -- ➜ stock=1",
        '',
        '-- 时刻 T2: Session B 查询库存',
        'SELECT stock FROM dishes',
        "  WHERE dish_id=1;          -- ➜ stock=1  (←也看到1!)",
        '',
        '-- 两个Session都认为库存充足 → 都下单 → 超卖!',
    ], 10)

    add_textbox(slide, 0.5, 3.8, 5.5, 1.5,
                '问题本质：SELECT 和 UPDATE 之间存在时间窗口。两个并发事务同时读到 stock=1，都判断"库存够"然后都下单。最终 2 个订单只扣了 1 份库存 — 这就是经典的 Lost Update 问题。',
                12, RED)

    # Right: FOR UPDATE solution
    add_textbox(slide, 6.5, 1.15, 6, 0.4, 'FOR UPDATE 排他锁解决方案', 16, GREEN, True)
    add_code_block(slide, 6.5, 1.6, 6.0, 3.0, [
        '-- 触发器 trg_check_dish_stock_before_order',
        '-- BEFORE INSERT ON order_items',
        'BEGIN',
        '  DECLARE v_stock INT;',
        '  -- ││ FOR UPDATE ││:',
        '  -- 对 dishes 该行加 X 锁(排他锁)',
        '  SELECT stock INTO v_stock',
        '    FROM dishes',
        '    WHERE dish_id = NEW.dish_id',
        '    FOR UPDATE;  -- ← 关键!',
        '',
        '  IF v_stock < NEW.quantity THEN',
        "    SIGNAL '45000'  -- 库存不足,拒绝",
        '  END IF;',
        'END',
    ], 10)

    add_textbox(slide, 6.5, 4.8, 6.0, 2.0,
                'FOR UPDATE 原理：\n'
                '[1] Session A 执行 SELECT...FOR UPDATE → 对 dish_id=1 加排他锁(X Lock)\n'
                '[2] Session B 执行 SELECT...FOR UPDATE → 阻塞等待 Session A 释放锁\n'
                '[3] Session A 完成事务(INSERT+COMMIT) → 释放锁\n'
                '[4] Session B 获得锁，重新读取 stock=0 → 库存不足 → 正确拒绝！',
                11, GREEN)

    # Bottom bar
    add_textbox(slide, 0.5, 6.6, 12, 0.5,
                '结论：FOR UPDATE 将"读-判断-写"三步原子化 — 这是数据库层解决并发超卖的唯一正确手段。应用层 synchronized 无法跨进程，乐观锁会频繁重试。',
                13, DARK_BLUE, True)

    page()

    # ===== SLIDE 18: Trigger 1-2: Stock Anti-Oversell =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '触发器 1-2：库存防超卖（BEFORE + AFTER 双触发器协同）', page())

    # Trigger 1
    add_textbox(slide, 0.3, 1.15, 6, 0.35, 'trg_check_dish_stock_before_order (BEFORE INSERT)', 14, DARK_BLUE, True)
    add_code_block(slide, 0.3, 1.55, 6.2, 2.3, [
        'CREATE TRIGGER trg_check_dish_stock_before_order',
        'BEFORE INSERT ON order_items',
        'FOR EACH ROW',
        'BEGIN',
        '  DECLARE v_stock INT;',
        '  DECLARE v_status TINYINT;',
        '  -- FOR UPDATE锁定菜品行，防止并发超卖',
        '  SELECT stock, status INTO v_stock, v_status',
        '    FROM dishes WHERE dish_id = NEW.dish_id',
        '    FOR UPDATE;  -- 排他锁',
        '  -- 检查1: 菜品是否已下架?',
        '  IF v_status = 0 THEN',
        "    SIGNAL SQLSTATE '45000'",
        "      SET MESSAGE_TEXT = '商品已下架!';",
        '  END IF;',
        '  -- 检查2: 库存是否充足?',
        '  IF v_stock < NEW.quantity THEN',
        "    SIGNAL SQLSTATE '45000'",
        "      SET MESSAGE_TEXT = '库存不足!';",
        '  END IF;',
        'END',
    ], 9)

    add_info_box(slide, 0.3, 4.0, 6.2, 1.6, '为什么需要 BEFORE INSERT 触发器?',
                 ['应用层检查库存有并发窗口 — 两个请求同时读到stock=1然后都下单',
                  'FOR UPDATE 对 dishes 行加 X 锁(排他锁)，其他事务的 FOR UPDATE 必须等待',
                  'MySQL 行级锁粒度精准到被查询的行，不影响其他菜品的并发下单',
                  '同时也检查下架状态(status=0)，防止买到已下架商品'])

    # Trigger 2
    add_textbox(slide, 6.8, 1.15, 6, 0.35, 'trg_reduce_dish_stock_after_order (AFTER INSERT)', 14, GREEN, True)
    add_code_block(slide, 6.8, 1.55, 5.8, 0.9, [
        'CREATE TRIGGER trg_reduce_dish_stock_after_order',
        'AFTER INSERT ON order_items',
        'FOR EACH ROW',
        'BEGIN',
        '  UPDATE dishes',
        '    SET stock = stock - NEW.quantity',
        '    WHERE dish_id = NEW.dish_id;',
        'END',
    ], 9)

    add_info_box(slide, 6.8, 2.6, 5.8, 1.4, '为什么 AFTER INSERT 而不是 BEFORE?',
                 ['只有BEFORE触发器通过所有校验（库存足+未下架），',
                  'AFTER触发器才会执行扣减。如果BEFORE拒绝，',
                  'AFTER不会执行 → 不会误扣库存',
                  '库存扣减是幂等操作，放在AFTER更安全'])

    add_info_box(slide, 6.8, 4.2, 5.8, 1.4, '触发时机与事务关系',
                 ['INSERT INTO order_items → trg_check(检查+加锁)',
                  '→ ✔✔ 通过 → trg_reduce(扣库存) → COMMIT',
                  '→ ✘ 拒绝 → ROLLBACK(整个事务回滚) → 不扣库存',
                  '两个触发器在同一个事务内，保证原子性'])

    page()

    # ===== SLIDE 19: Trigger 3-4: Rider Type Check =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '触发器 3-4：骑手类型约束（ENUM + 触发器双重保护）', page())

    add_textbox(slide, 0.5, 1.15, 12, 0.4,
                '问题背景：系统有 Stage1_Trunk(干线) 和 Stage2_Floor(楼栋) 两种骑手，如果错把楼栋骑手派去商家取餐，业务全乱。',
                13, DARK_TEXT)

    # Trigger 3
    add_textbox(slide, 0.3, 1.7, 6, 0.35, 'trg_check_rider_type_before_insert (BEFORE INSERT)', 14, DARK_BLUE, True)
    add_code_block(slide, 0.3, 2.1, 6.2, 1.8, [
        'CREATE TRIGGER trg_check_rider_type_before_insert',
        'BEFORE INSERT ON orders',
        'FOR EACH ROW',
        'BEGIN',
        '  IF NEW.stage1_rider_id IS NOT NULL THEN',
        '    IF NOT EXISTS (',
        '      SELECT 1 FROM riders',
        '      WHERE rider_id = NEW.stage1_rider_id',
        "        AND rider_type = 'Stage1_Trunk'",
        '    ) THEN',
        "      SIGNAL SQLSTATE '45000'",
        "        SET MESSAGE_TEXT = 'Stage1必须是Trunk类型!';",
        '    END IF;',
        '  END IF;',
        '  -- 同理检查 stage2_rider_id 必须是 Stage2_Floor',
        'END',
    ], 9)

    add_info_box(slide, 0.3, 4.1, 6.2, 2.5, 'INSERT 触发场景',
                 ['新订单插入时如果同时指定了骑手(INSERT INTO orders...stage1_rider_id=5)',
                  '触发器检查: rider_id=5 的 rider_type 是 Stage1_Trunk 吗?',
                  '✔ 是 → 允许写入',
                  '✘ 不是 → SIGNAL 拒绝 + 事务回滚',
                  '与 riders 表的 ENUM 约束形成双重保护:',
                  'ENUM 保证 rider_type 只有合法值，触发器保证用法正确'])

    # Trigger 4
    add_textbox(slide, 6.8, 1.7, 6, 0.35, 'trg_check_rider_type_before_update (BEFORE UPDATE)', 14, GREEN, True)
    add_code_block(slide, 6.8, 2.1, 5.8, 1.8, [
        'CREATE TRIGGER trg_check_rider_type_before_update',
        'BEFORE UPDATE ON orders',
        'FOR EACH ROW',
        'BEGIN',
        '  -- 仅在 stage1_rider_id 发生变化时检查',
        '  IF NEW.stage1_rider_id IS NOT NULL',
        '     AND (NEW.stage1_rider_id !=',
        '          OLD.stage1_rider_id',
        '          OR OLD.stage1_rider_id IS NULL)',
        '  THEN',
        '    -- 检查新指派的骑手类型是否正确',
        "    IF NOT EXISTS (...rider_type='Stage1_Trunk')",
        '    THEN SIGNAL...',
        '  END IF;',
        '  -- 同理检查 stage2_rider_id → Stage2_Floor',
        'END',
    ], 9)

    add_info_box(slide, 6.8, 4.1, 5.8, 2.5, 'UPDATE 触发场景与优化',
                 ['订单状态变更时重新指派骑手(UPDATE orders SET stage1_rider_id=8)',
                  '触发器检查: 仅在新旧值不同或旧值为NULL时触发检查',
                  '→ 避免每次UPDATE都重复校验（性能优化）',
                  'INSERT+UPDATE双触发器覆盖所有骑手指派路径',
                  '错误示例: 试图让Stage2_Floor骑手做Stage1配送 → 拒绝'])

    page()

    # ===== SLIDE 20: Trigger 5-6: Rider Status Auto-Management =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '触发器 5-6：骑手状态自动管理（Idle ↔ Delivering）', page())

    add_textbox(slide, 0.5, 1.15, 12, 0.4,
                '核心价值：应用层不需要手动 UPDATE riders SET status = ... — 所有状态转换由数据库触发器自动完成，杜绝遗漏和错误。',
                13, DARK_TEXT)

    # INSERT trigger
    add_textbox(slide, 0.3, 1.65, 6, 0.35, 'trg_rider_delivering_insert (AFTER INSERT)', 14, DARK_BLUE, True)
    add_code_block(slide, 0.3, 2.05, 6.2, 1.3, [
        'CREATE TRIGGER trg_rider_delivering_insert',
        'AFTER INSERT ON orders',
        'FOR EACH ROW',
        'BEGIN',
        '  -- 新订单已指派干线骑手 → 状态变为Delivering',
        "  IF NEW.stage1_rider_id IS NOT NULL",
        "     AND NEW.order_status IN ('Paid','Stage1_Assigned')",
        '  THEN',
        "    UPDATE riders SET status='Delivering'",
        '      WHERE rider_id = NEW.stage1_rider_id;',
        '  END IF;',
        '  -- 同理处理 stage2_rider_id',
        'END',
    ], 9)

    # UPDATE trigger - more complex
    add_textbox(slide, 6.8, 1.65, 6, 0.35, 'trg_rider_delivering_update (AFTER UPDATE)', 14, GREEN, True)
    add_code_block(slide, 6.8, 2.05, 5.8, 2.0, [
        'CREATE TRIGGER trg_rider_delivering_update',
        'AFTER UPDATE ON orders',
        'FOR EACH ROW',
        'BEGIN',
        '  -- [1] 新指派骑手 → Delivering',
        '  IF NEW.stage1_rider_id 发生变化 THEN',
        "    UPDATE riders SET status='Delivering'",
        '      WHERE rider_id = NEW.stage1_rider_id;',
        '  END IF;',
        '  -- [2] Stage1完成(Paid/Stage1_Assigned',
        '  --     → Arrived_At_Point) → 释放干线骑手',
        "  IF NEW.order_status='Arrived_At_Point'",
        "     AND OLD.order_status IN ('Paid','Stage1_Assigned')",
        '  THEN',
        "    UPDATE riders SET status='Idle'",
        '      WHERE rider_id = OLD.stage1_rider_id;',
        '  END IF;',
        '  -- [3] Stage2完成 → 释放楼栋骑手',
        '  -- [4] 取消 → 释放所有骑手',
        'END',
    ], 9)

    # Status transition table
    add_textbox(slide, 0.3, 3.6, 12, 0.4, '骑手状态自动转换表', 14, DARK_BLUE, True)
    add_table_simple(slide, 0.3, 4.1,
                     [1.8, 2.5, 2.5, 5.0],
                     ['触发事件', '订单状态变化', '骑手状态转换', '触发器逻辑'],
                     [
                         ['指派骑手', 'Paid → Stage1_Assigned', 'Idle → Delivering', 'AFTER INSERT/UPDATE 检测到新 stage1_rider_id'],
                         ['干线送达', 'Stage1_Assigned → Arrived_At_Point', 'Delivering → Idle', 'AFTER UPDATE: order_status=Arrived_At_Point 释放干线骑手'],
                         ['指派楼栋骑手', 'Arrived_At_Point → Stage2_Assigned', 'Idle → Delivering', 'AFTER UPDATE 检测到新 stage2_rider_id'],
                         ['最终送达', 'Stage2_Assigned → Completed', 'Delivering → Idle', 'AFTER UPDATE: order_status=Completed 释放楼栋骑手'],
                         ['订单取消', '任意 → Cancelled', 'Delivering → Idle', 'AFTER UPDATE: order_status=Cancelled 释放所有骑手'],
                     ], font_size=9)

    page()

    # ===== SLIDE 21: Trigger 7: Capacity Pre-Check (NEW!) =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '触发器 7（新增）：寄存点容量预检 — 防止骑手白跑', page())

    add_textbox(slide, 0.3, 1.15, 6, 0.35,
                'trg_check_pickup_point_capacity (BEFORE INSERT ON orders)', 14, DARK_BLUE, True)
    add_code_block(slide, 0.3, 1.55, 6.2, 1.5, [
        'CREATE TRIGGER trg_check_pickup_point_capacity',
        'BEFORE INSERT ON orders',
        'FOR EACH ROW',
        'BEGIN',
        '  DECLARE v_pt_current INT;',
        '  DECLARE v_pt_capacity INT;',
        '  -- FOR UPDATE锁住寄存点行，防并发容量击穿',
        '  SELECT current_packages, capacity',
        '    INTO v_pt_current, v_pt_capacity',
        '    FROM pickup_points',
        '    WHERE point_id = NEW.pickup_point_id',
        '    FOR UPDATE;  -- 排他锁',
        '  IF v_pt_current >= v_pt_capacity THEN',
        "    SIGNAL SQLSTATE '45000'",
        "      SET MESSAGE_TEXT = '该寄存点已满!'",
        '  END IF;',
        'END',
    ], 9)

    add_info_box(slide, 0.3, 3.2, 6.2, 1.5, '为什么需要这个新触发器?',
                 ['之前只有 chk_capacity 硬约束 — 在入库时拦截(等骑手送到了才发现柜满)',
                  '→ 骑手白跑一趟，已完成的工作被 ROLLBACK 抹掉',
                  '新触发器在下单时就检查容量 → 满了直接拒绝，用户选其他寄存点',
                  'FOR UPDATE 保证两个用户不会同时看到"还剩1格"同时下单'])

    # Right side: Two-layer protection diagram
    add_textbox(slide, 6.8, 1.15, 6, 0.35, '爆仓防护：两层体系', 16, DARK_BLUE, True)

    # Layer 1
    card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.8), Inches(1.5))
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    card1.line.color.rgb = GREEN
    card1.line.width = Pt(2)
    add_textbox(slide, 7.0, 1.8, 5.4, 0.3, '第一层：下单容量预检（主动防护）', 14, GREEN, True)
    add_multiline_textbox(slide, 7.0, 2.2, 5.4, 0.9, [
        'sp_create_order: FOR UPDATE查容量 → 满→拒绝("请选邻近寄存点")',
        'trg_check_pickup_point_capacity: BEFORE INSERT触发',
        '→ 用户下单时就被引导分流，不会出现骑手送到后发现柜满'], 11, DARK_TEXT)

    # Layer 2
    card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(3.5), Inches(5.8), Inches(1.5))
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    card2.line.color.rgb = RED
    card2.line.width = Pt(2)
    add_textbox(slide, 7.0, 3.6, 5.4, 0.3, '第二层：物理硬约束（被动兜底）', 14, RED, True)
    add_multiline_textbox(slide, 7.0, 4.0, 5.4, 0.9, [
        'chk_capacity CHECK (current_packages <= capacity)',
        'sp_arrive_at_pickup_point: UPDATE current_packages+1',
        '→ 如果+1后>capacity → ROLLBACK(安全网, 正常情况走不到这步)'], 11, DARK_TEXT)

    # Bottom: summary
    add_textbox(slide, 0.3, 5.2, 13, 0.8,
                '两层关系：第一层在下单时主动拦(用户体验好)，第二层在入库时兜底(极端并发+绕过SP的直接INSERT)。\n'
                'FOR UPDATE 行级锁确保两层各自的容量检查都是原子性的，不会出现"同时看到最后一格"的并发击穿。',
                12, DARK_BLUE)

    add_info_box(slide, 0.3, 6.1, 12.5, 0.9, '数据库层实际拦截了什么?',
                 ['整个防护链: trg_check(下单时容量预检) → FOR UPDATE(并发保护) → chk_capacity(入库时硬约束)',
                  '拦截点1: INSERT INTO orders → 容量满→拒绝 → 用户换寄存点',
                  '拦截点2: sp_arrive_at_pickup_point → current_packages+1>capacity → ROLLBACK → 包裹回退到Stage1_Assigned'])

    page()

    # ===== SLIDE 22: All 7 Triggers Summary =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '七重触发器总览', page())

    add_table_simple(slide, 0.3, 1.3,
                     [0.4, 3.0, 1.4, 1.3, 1.3, 5.0],
                     ['#', '触发器名称', '触发时机', '事件', '表', '功能描述'],
                     [
                         ['1', 'trg_check_dish_stock_before_order', 'BEFORE', 'INSERT', 'order_items', 'FOR UPDATE锁定菜品行 → 检查库存充足+未下架 → 不足则SIGNAL拒绝'],
                         ['2', 'trg_reduce_dish_stock_after_order', 'AFTER', 'INSERT', 'order_items', '自动扣减菜品库存: stock = stock - quantity'],
                         ['3', 'trg_check_rider_type_before_insert', 'BEFORE', 'INSERT', 'orders', '校验stage1_rider必须是Stage1_Trunk, stage2必须是Stage2_Floor'],
                         ['4', 'trg_check_rider_type_before_update', 'BEFORE', 'UPDATE', 'orders', '同#3，在骑手变更时触发(仅新旧值不同时校验)'],
                         ['5', 'trg_rider_delivering_insert', 'AFTER', 'INSERT', 'orders', '新订单指派骑手 → 自动设置 ridER.status=Delivering'],
                         ['6', 'trg_rider_delivering_update', 'AFTER', 'UPDATE', 'orders', '状态驱动: 指派→Delivering, 完成→Idle, 取消→Idle'],
                         ['7', 'trg_check_pickup_point_capacity', 'BEFORE', 'INSERT', 'orders', 'FOR UPDATE查寄存点容量 → 满则SIGNAL拒绝 → 用户选其他点'],
                     ], font_size=8)

    add_textbox(slide, 0.3, 5.6, 12, 0.8,
                '分类统计: 库存管理 2个(触发器1-2) | 骑手类型约束 2个(触发器3-4) | 骑手状态自动管理 2个(触发器5-6) | 容量预检 1个(触发器7)\n'
                '技术栈: FOR UPDATE行级锁(触发器1,7) | SIGNAL异常处理(触发器1,3,4,7) | NEW/OLD状态对比(触发器4,6) | ENUM类型+NOT EXISTS(触发器3,4)',
                11, GRAY_TEXT)

    page()

    # ===== SLIDE 23: SECTION 05 - Stored Procedures =====
    add_section_header(prs, '05', '存储过程与事务管理',
                       '4个核心SP · 原子事务 · ROLLBACK回滚 · 游标 · 爆仓处理')
    page()

    # ===== SLIDE 24: sp_create_order =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, 'sp_create_order：原子性下单（含余额检查 + 容量预检 + 扣款）', page())

    add_code_block(slide, 0.3, 1.2, 7.5, 5.0, [
        'CREATE PROCEDURE sp_create_order(',
        '  IN p_user_id INT, IN p_merchant_id INT,',
        '  IN p_point_id INT, IN p_dish_id INT,',
        '  IN p_quantity INT, OUT o_order_id INT)',
        'BEGIN',
        '  DECLARE EXIT HANDLER FOR SQLEXCEPTION',
        '  BEGIN  ROLLBACK;  RESIGNAL;  END;',
        '  START TRANSACTION;  -- 事务开始',
        '',
        '  -- [1] 获取菜品价格',
        '  SELECT price INTO v_dish_price FROM dishes',
        '    WHERE dish_id = p_dish_id;',
        '',
        '  -- [2] 检查钱包余额',
        '  SELECT balance INTO v_balance FROM users',
        '    WHERE user_id = p_user_id;',
        '  IF v_balance < total THEN SIGNAL... END IF;',
        '',
        '  -- [3] 容量预检 (FOR UPDATE防并发)',
        '  SELECT current_packages, capacity',
        '    INTO v_pt_current, v_pt_capacity',
        '    FROM pickup_points',
        '    WHERE point_id = p_point_id FOR UPDATE;',
        '  IF v_pt_current >= v_pt_capacity THEN',
        "    SIGNAL '该寄存点已满!';",
        '  END IF;',
        '',
        '  -- [4] 扣款 + [5] 插入订单 + [6] 插入明细',
        '  UPDATE users SET balance = balance - total;',
        '  INSERT INTO orders (...) VALUES (...);',
        '  INSERT INTO order_items (...) VALUES (...);',
        '  -- 明细INSERT触发 trg_check(锁库存) + trg_reduce(扣库存)',
        '  COMMIT;  -- 原子提交',
        'END',
    ], 8)

    add_textbox(slide, 8.2, 1.2, 4.5, 0.3, '执行流程与触发器链', 14, DARK_BLUE, True)
    flow = [
        '1. START TRANSACTION',
        '2. 查菜品价格 → 算总金额',
        '3. 查钱包余额 → 不足→SIGNAL→ROLLBACK',
        '4. FOR UPDATE查寄存点容量 → 满→SIGNAL→ROLLBACK',
        '5. UPDATE扣钱包',
        '6. INSERT INTO orders',
        '   → 触发 trg_check_rider_type (骑手类型校验)',
        '   → 触发 trg_check_pickup_point_capacity (容量预检)',
        '7. INSERT INTO order_items',
        '   → 触发 trg_check_dish_stock (FOR UPDATE库存检查)',
        '   → 触发 trg_reduce_dish_stock (库存扣减)',
        '8. COMMIT (所有步骤全部成功)',
        '任一环节失败 → EXIT HANDLER捕获',
        '→ ROLLBACK → RESIGNAL → 上游收到错误',
    ]
    add_multiline_textbox(slide, 8.2, 1.6, 4.5, 5.0, flow, 10, DARK_TEXT, line_spacing=1.3)

    page()

    # ===== SLIDE 25: sp_arrive + sp_stage2_deliver =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, 'sp_arrive_at_pickup_point + sp_stage2_deliver：两段配送引擎', page())

    # Left: Stage1 arrival
    add_textbox(slide, 0.3, 1.15, 6, 0.35, 'sp_arrive_at_pickup_point — 干线骑手送达入库', 14, DARK_BLUE, True)
    add_code_block(slide, 0.3, 1.55, 6.0, 2.8, [
        'CREATE PROCEDURE sp_arrive_at_pickup_point(',
        '  IN p_order_id INT)',
        'BEGIN',
        '  DECLARE EXIT HANDLER FOR SQLEXCEPTION',
        '  BEGIN ROLLBACK; RESIGNAL; END;',
        '  START TRANSACTION;',
        '',
        '  -- Step 1: 状态跳转(可逆)',
        '  UPDATE orders',
        "    SET order_status = 'Arrived_At_Point',",
        '        stage1_completed_at = CURRENT_TIMESTAMP',
        '    WHERE order_id = p_order_id;',
        '',
        '  -- Step 2: 寄存点计数+1(触发chk_capacity)',
        '  UPDATE pickup_points',
        '    SET current_packages = current_packages + 1',
        '    WHERE point_id = v_point_id;',
        '  -- ↑ chk_capacity CHECK(current<=capacity)',
        '  -- 如果溢出 → 异常 → ROLLBACK',
        '  -- Step1的状态变更有被回滚',
        '',
        '  -- 触发器自动: stage1_rider → Idle',
        '  COMMIT;',
        'END',
    ], 8)

    add_info_box(slide, 0.3, 4.55, 6.0, 2.0, '爆仓时的行为',
                 ['若 current_packages+1 > capacity:',
                  'chk_capacity CHECK → 异常 → ROLLBACK',
                  '订单状态回退(Arrived_At_Point→Stage1_Assigned)',
                  'stage1_completed_at 被回滚(抹掉了干线骑手的工作)',
                  '→ 这是两段式配送的不足: 骑手已完成但系统不认',
                  '改进方向: 引入"已到达但等待空位"中间状态'])

    # Right: Stage2 delivery
    add_textbox(slide, 6.8, 1.15, 6, 0.35, 'sp_stage2_deliver — 楼栋骑手送达完成', 14, GREEN, True)
    add_code_block(slide, 6.8, 1.55, 5.8, 1.8, [
        'CREATE PROCEDURE sp_stage2_deliver(',
        '  IN p_order_id INT)',
        'BEGIN',
        '  DECLARE EXIT HANDLER FOR SQLEXCEPTION',
        '  BEGIN ROLLBACK; RESIGNAL; END;',
        '  START TRANSACTION;',
        '',
        '  -- Step 1: 标记订单完成',
        '  UPDATE orders',
        "    SET order_status = 'Completed',",
        '        stage2_completed_at = CURRENT_TIMESTAMP',
        '    WHERE order_id = p_order_id;',
        '',
        '  -- Step 2: 寄存点包裹数-1',
        '  UPDATE pickup_points',
        '    SET current_packages = current_packages - 1',
        '    WHERE point_id = v_point_id',
        '      AND current_packages > 0;  -- 安全保护',
        '',
        '  -- 触发器自动: stage2_rider → Idle',
        '  COMMIT;',
        'END',
    ], 8)

    add_info_box(slide, 6.8, 3.55, 5.8, 3.0, 'sp_cancel_order 取消流程',
                 ['1. 查订单状态(Completed/Cancelled不能再次取消)',
                  '2. UPDATE users: balance + total(退款)',
                  '3. 游标遍历 order_items → 恢复每道菜库存',
                  '4. 如果已入库 → current_packages - 1',
                  '5. UPDATE orders SET status=Cancelled',
                  '→ 触发器自动释放所有骑手(Idle)',
                  '所有步骤在同一事务内，失败→ROLLBACK'])

    page()

    # ===== SLIDE 26: Transaction & Rollback =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '事务回滚机制与爆仓处理流程', page())

    # ROLLBACK diagram with text
    add_textbox(slide, 0.5, 1.2, 5.5, 0.4, 'EXIT HANDLER 异常捕获与回滚机制', 16, DARK_BLUE, True)
    add_code_block(slide, 0.5, 1.7, 5.5, 1.5, [
        '-- 每个存储过程都有的标准模板',
        'DECLARE EXIT HANDLER FOR SQLEXCEPTION',
        'BEGIN',
        '  ROLLBACK;  -- 撤销事务内所有操作',
        '  RESIGNAL;  -- 向上层抛出具体错误信息',
        'END;',
        '',
        'START TRANSACTION;',
        '  ... 业务操作 ...',
        '  -- 任意步骤失败 → 自动跳到HANDLER',
        'COMMIT;  -- 全部成功才提交',
    ], 10)

    # Right: Flow diagram
    add_textbox(slide, 6.5, 1.2, 6, 0.4, '爆仓处理完整流程', 16, RED, True)
    flow_lines = [
        '[1] 用户下单 → sp_create_order',
        '    → 容量预检(FOR UPDATE) → 满→拒绝("请选其他寄存点")',
        '    → 有空 → 下单成功 → 订单状态=Paid',
        '',
        '[2] 干线骑手取件 → 订单状态=Stage1_Assigned',
        '',
        '[3] 干线骑手送达 → sp_arrive_at_pickup_point',
        '    → UPDATE current_packages+1',
        '    → chk_capacity CHECK(current<=capacity)',
        '    → 如果溢出(极端并发): ROLLBACK, 订单回Stage1_Assigned',
        '    → 如果正常: 入库成功, 订单=Arrived_At_Point',
        '',
        '[4] 楼栋骑手取件 → 订单=Stage2_Assigned',
        '',
        '[5] 楼栋骑手送达 → sp_stage2_deliver',
        '    → current_packages-1 → 释放一个格子',
        '    → 订单=Completed → 全链路完成',
    ]
    add_multiline_textbox(slide, 6.5, 1.7, 6.0, 5.0, flow_lines, 10, DARK_TEXT, line_spacing=1.2)

    # Bottom summary
    add_textbox(slide, 0.5, 3.5, 5.5, 3.5,
                '事务四大特性(ACID)在本系统中的体现:\n\n'
                'A-原子性: 所有SP用START TRANSACTION+COMMIT/ROLLBACK，\n'
                '  扣款+插订单+插明细要么全成功要么全回滚\n\n'
                'C-一致性: CHECK约束(chk_capacity, balance>=0, stock>=0)\n'
                '  保证数据库永远处于合法状态\n\n'
                'I-隔离性: FOR UPDATE行级锁，并发事务之间互不干扰\n'
                '  两个用户同时买最后一份库存 → 一个成功一个等待\n\n'
                'D-持久性: InnoDB引擎，COMMIT后数据写入redo log\n'
                '  系统崩溃后可从日志恢复已提交的事务',
                11, DARK_TEXT)

    page()

    # ===== SLIDE 27: SECTION 06 - Views =====
    add_section_header(prs, '06', '视图与数据分析',
                       '饱和度预警视图 · RANK() OVER窗口函数 · 商户销售排行')
    page()

    # ===== SLIDE 28: Views Detailed =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '2个业务视图：驱动大屏实时数据展示', page())

    # View 1
    add_textbox(slide, 0.3, 1.15, 6, 0.35, 'vw_pickup_point_analytics — 寄存点饱和度预警视图', 14, DARK_BLUE, True)
    add_code_block(slide, 0.3, 1.55, 6.0, 2.0, [
        'CREATE VIEW vw_pickup_point_analytics AS',
        'SELECT',
        '  p.point_id,',
        '  p.point_name,                  -- 寄存点名称',
        '  p.capacity AS max_capacity,    -- 最大容量',
        '  p.current_packages,            -- 当前包裹数',
        '  ROUND(p.current_packages/p.capacity*100,2)',
        '    AS saturation_pct,           -- 饱和度百分比',
        '  COALESCE(sub.backlog_count,0)  -- 滞留件数',
        '    AS backlog_count',
        'FROM pickup_points p',
        'LEFT JOIN (',
        '  SELECT pickup_point_id, COUNT(*) AS backlog_count',
        '  FROM orders',
        "  WHERE order_status IN ('Arrived_At_Point',",
        "    'Stage2_Assigned')",
        '  GROUP BY pickup_point_id',
        ') sub ON p.point_id = sub.pickup_point_id;',
    ], 8)

    add_info_box(slide, 0.3, 3.75, 6.0, 2.8, '大屏爆仓预警如何使用这个视图',
                 ['大屏每30秒查询: SELECT * FROM vw_pickup_point_analytics',
                  'saturation_pct > 80% → 黄色预警 (接近爆仓)',
                  'saturation_pct = 100% → 红色爆仓 (已满，新订单拒绝)',
                  'backlog_count 与 current_packages 理论一致(chk_capacity保证)',
                  '新索引 idx_orders_point_status 加速子查询LEFT JOIN',
                  '驱动大屏"爆仓预警"Tab的柱状图和数据卡片'])

    # View 2
    add_textbox(slide, 6.8, 1.15, 6, 0.35, 'vw_merchant_sales_rank — 商户销售排行视图', 14, GREEN, True)
    add_code_block(slide, 6.8, 1.55, 5.8, 1.5, [
        'CREATE VIEW vw_merchant_sales_rank AS',
        'SELECT',
        '  m.merchant_id,',
        '  m.merchant_name,               -- 商家名称',
        '  COUNT(DISTINCT o.order_id)',
        '    AS total_orders,             -- 订单总数',
        '  IFNULL(SUM(o.total_amount),0)',
        '    AS total_sales,              -- 销售总额',
        '  RANK() OVER (',
        '    ORDER BY IFNULL(SUM(o.total_amount),0) DESC',
        '  ) AS sales_rank                -- 销售额排名',
        'FROM merchants m',
        'LEFT JOIN orders o',
        "  ON m.merchant_id = o.merchant_id",
        "  AND o.order_status = 'Completed'",
        'GROUP BY m.merchant_id;',
    ], 8)

    add_info_box(slide, 6.8, 3.25, 5.8, 3.3, 'RANK() OVER 窗口函数解析',
                 ['RANK() OVER (ORDER BY SUM(amount) DESC)',
                  '',
                  '传统做法: 子查询+变量，SQL复杂且易出错',
                  '窗口函数: 一行SQL完成排名，数据库引擎优化执行',
                  '',
                  'RANK() vs ROW_NUMBER() vs DENSE_RANK():',
                  '• RANK(): 同销售额排名相同，下一个跳过(1,2,2,4)',
                  '• DENSE_RANK(): 同销售额排名相同，连续(1,2,2,3)',
                  '• ROW_NUMBER(): 无并列，随机排序(1,2,3,4)',
                  '',
                  '仅统计Completed订单 → 排除未完成交易的干扰'])

    page()

    # ===== SLIDE 29: SECTION 07 - Dashboard =====
    add_section_header(prs, '07', '系统大屏与可视化展示',
                       'Flask + ECharts 实时监控大屏 · 3个功能Tab · 10张截图')
    page()

    # ===== SLIDE 30: Dashboard Overview =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '大屏总览：实时数据监控 Dashboard', page())
    img = insert_image_safe(slide, os.path.join(IMG_DIR, '01_dashboard_overview.png'),
                            0.3, 1.2, 12.5, 5.8)
    if not img:
        add_textbox(slide, 1, 3, 10, 1, '[截图未找到，请启动 Flask 后截取大屏全貌]', 16, RED)
    page()

    # ===== SLIDE 31: Order Section + KPI =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '订单管理模块与 KPI 指标卡', page())
    insert_image_safe(slide, os.path.join(IMG_DIR, '02_order_section.png'),
                      0.3, 1.1, 6.2, 3.0)
    insert_image_safe(slide, os.path.join(IMG_DIR, '04_kpi_top.png'),
                      6.8, 1.1, 6.0, 3.0)
    add_textbox(slide, 0.3, 4.3, 12.5, 1.5,
                '左侧：订单管理 Tab — 按状态筛选(Paid/Stage1_Assigned/Arrived_At_Point/Stage2_Assigned/Completed/Cancelled)，实时数据表格\n'
                '右侧：KPI 指标卡 — 今日订单数、今日营收、活跃骑手数、活跃商家数、爆仓寄存点数，30秒自动刷新',
                12, DARK_TEXT)
    page()

    # ===== SLIDE 32: Overflow Warning =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '爆仓预警：饱和度实时监控', page())
    insert_image_safe(slide, os.path.join(IMG_DIR, '10_overflow_dashboard.png'),
                      0.3, 1.1, 7.5, 3.0)
    insert_image_safe(slide, os.path.join(IMG_DIR, '11_overflow_pickup.png'),
                      8.0, 1.1, 4.8, 3.0)
    add_textbox(slide, 0.3, 4.3, 12.5, 2.0,
                '爆仓预警机制:\n'
                '[1] chk_capacity CHECK约束 — 数据库物理层硬约束，current_packages绝对不能超过capacity\n'
                '[2] trg_check_pickup_point_capacity — 下单前FOR UPDATE查容量，满了直接拒绝\n'
                '[3] vw_pickup_point_analytics — 实时计算saturation_pct，>80%黄色预警，=100%红色爆仓\n'
                '[4] API /api/today_stats — 统计 overflow_points(饱和度>80%的寄存点数量)\n'
                '当前模拟数据: 3期=100%(爆仓), 6期=100%(爆仓), 其余正常 → 验证两层防护体系有效',
                11, DARK_TEXT)
    page()

    # ===== SLIDE 33: Other Management Pages =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '管理页面：寄存点 / 商家 / 学生 / 骑手 / 菜品', page())

    screenshots_data = [
        ('03_pickup_points.png', '寄存点管理', 0.3, 1.1, 4.0, 2.6),
        ('05_merchants.png', '商家管理', 4.5, 1.1, 4.0, 2.6),
        ('06_students.png', '学生管理', 8.8, 1.1, 4.0, 2.6),
        ('07_riders.png', '骑手管理', 0.3, 3.9, 4.0, 2.6),
        ('08_dishes.png', '菜品管理', 4.5, 3.9, 4.0, 2.6),
    ]
    for fname, label, x, y, w, h in screenshots_data:
        insert_image_safe(slide, os.path.join(IMG_DIR, fname), x, y, w, h)
        add_textbox(slide, x, y + h + 0.05, w, 0.3, label, 11, GRAY_TEXT, align=PP_ALIGN.CENTER)

    page()

    # ===== SLIDE 34: AI Text-to-SQL =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, 'AI Text-to-SQL：自然语言查询数据库', page())

    add_textbox(slide, 0.5, 1.2, 12, 0.4,
                '集成 DeepSeek API，用户在搜索框输入自然语言 → AI 生成 SQL → 执行并返回结果', 14, DARK_BLUE, True)

    add_code_block(slide, 0.5, 1.8, 7.0, 2.5, [
        '# Python 提示词工程 (app.py)',
        "system_prompt = '''",
        '你是一个MySQL查询助手。数据库有以下表:',
        '- users(user_id, username, phone, dorm_building, balance)',
        '- merchants(merchant_id, merchant_name, address, rating)',
        '- dishes(dish_id, merchant_id, dish_name, price, stock, status)',
        '- pickup_points(point_id, point_name, location, capacity, current_packages)',
        '- riders(rider_id, rider_name, rider_type, status)',
        '- orders(order_id, user_id, merchant_id, pickup_point_id,',
        '  total_amount, order_status, stage1_rider_id, stage2_rider_id)',
        '- order_items(item_id, order_id, dish_id, quantity, price_at_order)',
        '',
        '请根据用户自然语言输入，生成单个MySQL SELECT查询语句。',
        "不要使用INSERT/UPDATE/DELETE/DROP。只返回纯SQL，不要解释。",
        "'''",
    ], 8)

    add_info_box(slide, 8.0, 1.8, 4.8, 2.0, 'AI 查询示例',
                 ['输入: "今天卖了多少钱?"',
                  '→ SQL: SELECT SUM(total_amount) FROM orders',
                  '  WHERE DATE(created_at)=CURDATE()',
                  '  AND order_status=\'Completed\'',
                  '',
                  '输入: "哪个寄存点快满了?"',
                  '→ SQL: SELECT point_name, saturation_pct',
                  '  FROM vw_pickup_point_analytics',
                  '  ORDER BY saturation_pct DESC'])

    add_info_box(slide, 8.0, 4.0, 4.8, 2.5, 'Schema注入策略',
                 ['提示词中包含完整数据库Schema',
                  '让AI理解表结构、字段含义、ENUM值',
                  '限制只生成SELECT(安全考虑)',
                  '前端展示SQL和执行结果(透明可信)',
                  'DeepSeek模型: deepseek-chat',
                  '成本极低: 单次查询<0.01元人民币'])

    add_textbox(slide, 0.5, 4.6, 7.0, 2.0,
                '技术亮点:\n\n'
                '[1] Schema注入 — 将完整数据库结构放入系统提示词，AI能准确理解表关系和字段含义\n'
                '[2] 安全限制 — 只允许SELECT语句，拒绝任何写操作，防止AI幻觉导致数据损坏\n'
                '[3] 透明展示 — 前端同时展示生成的SQL和查询结果，用户可以看到AI"在想什么"\n'
                '[4] 视图复用 — AI可以直接查询 vw_pickup_point_analytics 和 vw_merchant_sales_rank，无需重复JOIN',
                11, DARK_TEXT)

    page()

    # ===== SLIDE 35: SECTION 08 - Summary =====
    add_section_header(prs, '08', '创新点总结与成果交付',
                       '技术创新 · 数据库设计亮点 · 项目交付物 · 未来展望')
    page()

    # ===== SLIDE 36: Innovation Points =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '七大创新点', page())

    innovations = [
        ('FOR UPDATE 行级锁防超卖', '触发器内嵌 SELECT...FOR UPDATE，高并发场景下原子化"读-判断-写"，杜绝超卖。这是数据库课程的核心考核点。', DARK_BLUE),
        ('七重触发器防护体系', '库存防超卖(2个) + 骑手类型约束(2个) + 骑手状态自动管理(2个) + 容量预检(1个)。数据库层全自动，应用层零负担。', DARK_BLUE),
        ('两段式配送状态机', '6态ENUM + 双骑手绑定 + 3时间戳审计。覆盖校园外卖全流程从下单到签收，每个状态转换都有明确触发条件。', GREEN),
        ('chk_capacity 物理硬约束', 'CHECK(current_packages <= capacity) 确保数据库层物理容积不超限。配合容量预检形成两层防护。', GREEN),
        ('RANK() OVER 窗口函数', '商户销售排行视图，一行SQL替代复杂的子查询+变量漂移。数据库引擎优化执行，比应用层排序更快更准确。', ORANGE),
        ('AI Text-to-SQL', 'DeepSeek LLM + Schema注入提示词，自然语言→SQL→结果。前端透明展示SQL，安全限制只读。', ORANGE),
        ('Cloudflare Tunnel 公网部署', '免费隧道技术将 localhost:5000 映射到公网域名，答辩时评委可直接手机扫码访问实时大屏。', RGBColor(0x7B, 0x1F, 0xA2)),
    ]
    for i, (title, desc, clr) in enumerate(innovations):
        y = 1.1 + i * 0.85
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.5), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = clr
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER

        add_textbox(slide, 1.2, y, 3.2, 0.5, title, 14, clr, True)
        add_textbox(slide, 4.5, y + 0.05, 8.3, 0.65, desc, 11, DARK_TEXT)

    page()

    # ===== SLIDE 37: Deliverables =====
    slide = prs.slides.add_slide(BLANK)
    add_title_bar(slide, '项目交付成果', page())

    deliverables = [
        ('数据库设计', ['campus_delivery_db.sql (完整DDL+DML)', '7表 4索引 7触发器 4存储过程 2视图', '严格3NF + 完备外键级联 + CHECK约束']),
        ('模拟数据', ['generate_mock_data.py (5000条订单)', '容量感知生成器(backlog<=capacity)', '生成数据已通过所有触发器校验']),
        ('Flask大屏', ['app.py (REST API + HTML渲染)', '3个功能Tab + ECharts可视化', '30秒自动刷新 + AI Text-to-SQL']),
        ('文档', ['校园外卖两段式配送系统_实践报告.docx', 'README.md (项目全中文化)', '本文档 PPT 答辩汇报']),
        ('部署', ['Cloudflare Tunnel 公网访问', 'proxy.py 代理管理脚本', '.env 环境配置管理']),
    ]
    for i, (title, items) in enumerate(deliverables):
        x = 0.3 + i * 2.55
        add_card(slide, x, 1.2, 2.35, 4.5, title, items,
                 [DARK_BLUE, GREEN, ORANGE, RGBColor(0x7B, 0x1F, 0xA2), RED][i])

    add_textbox(slide, 0.5, 6.0, 12, 0.5,
                'Git 仓库: 校园外卖两段式配送数据库系统  |  所有代码已提交推送  |  分支: master',
                12, GRAY_TEXT)

    page()

    # ===== SLIDE 38: Future + Thanks =====
    slide = prs.slides.add_slide(BLANK)
    # Full dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

    add_textbox(slide, 1.0, 1.5, 11, 1.0, '未来展望与改进方向', 32, WHITE, True)

    future = [
        '[1] 柜满排队机制 — 引入"已到达但等待空位"中间状态，不抹掉干线骑手已完成的工作',
        '[2] 智能容量调度 — 下单时自动推荐邻近有空位的寄存点，用户一键切换',
        '[3] 配送路线优化 — 楼栋骑手多单拼单路径规划，一次取5-8件集中配送',
        '[4] 实时骑手GPS — 接入地图API，用户实时追踪包裹位置',
        '[5] 全链路压测 — sysbench/ab工具测试QPS上限，验证FOR UPDATE性能瓶颈',
        '[6] 微服务拆分 — 订单服务/骑手服务/寄存点服务独立部署，消息队列解耦',
    ]
    add_multiline_textbox(slide, 1.0, 2.8, 11, 3.0, future, 14, RGBColor(0xCC, 0xCC, 0xCC), line_spacing=1.8)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(5.8), Inches(7), Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = DARK_BLUE
    div.line.fill.background()

    add_textbox(slide, 1.0, 6.2, 11, 0.8, '感谢各位老师聆听！\n欢迎扫码访问实时大屏  |  GitHub: sou1maker', 16, GRAY_TEXT, align=PP_ALIGN.CENTER)

    page()

    # ===== SAVE =====
    output_path = os.path.join(os.path.dirname(__file__),
                               '校园外卖两段式配送系统_答辩汇报.pptx')
    prs.save(output_path)
    print(f'\nPPT saved: {output_path}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    build_ppt()
